import csv
import os
import logging
from typing import Optional
from PIL import Image
import fitz  # PyMuPDF
import pdfplumber
import pytesseract
import openpyxl
import docx  # python-docx
from pptx import Presentation  # python-pptx

from fastapi import UploadFile
from sqlalchemy.orm import Session
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.core.database import SessionLocal
from app.schemas.document import DocumentOut
from app.models.document import Document, DocumentStatus
from app.repository import document_repo
from app.repository.conversation_repo import get_conversation

from app.service.rag_clients import (
    embedding_manager,
    vector_store,
)

from app.service.tools.image_tool import (
    generate_image_caption,
    ImageCaptionQuotaExceededError,
)


logger = logging.getLogger(__name__)


UPLOAD_DIR = "Uploads"

os.makedirs(
    UPLOAD_DIR,
    exist_ok=True,
)


# ============================================================
# PDF / PPTX IMAGE CAPTION LIMIT
# ============================================================

MAX_PDF_IMAGES_TO_CAPTION = int(
    os.getenv(
        "MAX_PDF_IMAGES_TO_CAPTION",
        "15",
    )
)


# ============================================================
# PER-PAGE RENDERED FALLBACK THRESHOLD (§12)
#
# A PDF page whose extracted text is shorter than this (in
# characters) is treated as "mostly visual" — a diagram, chart,
# or screenshot that isn't exposed as a normal embedded image
# object. Such pages are rendered and processed as images IN
# ADDITION to whatever embedded images/text that page already
# has, as long as that specific page had zero embedded images
# extracted (to avoid duplicating a page that's already covered).
# ============================================================

PDF_PAGE_TEXT_FALLBACK_THRESHOLD = int(
    os.getenv(
        "PDF_PAGE_TEXT_FALLBACK_THRESHOLD",
        "40",
    )
)


# ============================================================
# SUPPORTED FILE TYPES
# ============================================================

IMAGE_MIME_TYPES = {
    "image/jpeg",
    "image/png",
    "image/webp",
    "image/gif",
}


IMAGE_EXTENSIONS = {
    ".jpg",
    ".jpeg",
    ".png",
    ".webp",
    ".gif",
}


TEXT_EXTENSIONS = {
    ".md",
    ".txt",
}


CSV_EXTENSIONS = {
    ".csv",
}


EXCEL_EXTENSIONS = {
    ".xlsx",
    ".xlsm",
}


DOCX_EXTENSIONS = {
    ".docx",
}


DOCX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}


PPTX_EXTENSIONS = {
    ".pptx",
}


PPTX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


# ============================================================
# GCS PATH
# ============================================================


def generate_gcs_path(
    user_id: str,
    document_id: str,
    original_filename: str,
) -> str:

    safe_filename = (
        original_filename
        .replace("/", "_")
        .strip()
    )

    return (
        f"users/{user_id}/documents/"
        f"{document_id}/{safe_filename}"
    )


def get_local_file_path(doc: Document) -> str:
    """
    The on-disk path convention used everywhere a document's raw
    file is saved/read (upload, PDF/PPTX image extraction,
    download, live captioning). Centralised here so every caller
    stays in sync instead of re-building the same string
    independently.
    """

    return os.path.join(
        UPLOAD_DIR,
        f"{doc.id}_{doc.file_name}",
    )


# ============================================================
# CREATE FOLDER
# ============================================================


def create_folder_service(
    db: Session,
    file_name: str,
    parent_id: Optional[str],
    user_id: str,
) -> Optional[Document]:

    if not parent_id:
        parent_id = None

    if parent_id:

        parent = document_repo.get_owned_folder_by_id(
            db,
            parent_id,
            user_id,
        )

        if not parent:
            return None

    return document_repo.create_folder(
        db,
        file_name,
        parent_id,
        user_id,
    )


# ============================================================
# VALIDATE DOCUMENT CONVERSATION
# ============================================================


def _validate_document_conversation(
    db: Session,
    conversation_id: Optional[str],
    user_id: str,
) -> Optional[str]:
    """
    Validate the conversation that will own the uploaded document.

    conversation_id=None:
        Global document.

    conversation_id provided:
        Conversation must belong to current user.
    """

    if not conversation_id:
        return None

    conversation = get_conversation(
        db=db,
        conversation_id=str(conversation_id),
        user_id=str(user_id),
    )

    if not conversation:
        raise PermissionError(
            "Conversation not found or does not belong "
            "to the current user"
        )

    return str(conversation.id)


# ============================================================
# UPLOAD DOCUMENT
# ============================================================


def upload_document_service(
    db: Session,
    file: UploadFile,
    parent_id: Optional[str],
    conversation_id: Optional[str],
    user_id: str,
) -> Optional[Document]:

    logger.info(
        "UPLOAD DEBUG | file=%s | parent_id=%r | "
        "conversation_id=%r | user_id=%r",
        file.filename,
        parent_id,
        conversation_id,
        user_id,
    )

    if not parent_id:
        parent_id = None

    if not conversation_id:
        conversation_id = None

    # --------------------------------------------------------
    # GLOBAL FOLDER VALIDATION
    # --------------------------------------------------------

    if parent_id:

        parent = document_repo.get_owned_folder_by_id(
            db,
            parent_id,
            user_id,
        )

        if not parent:
            return None

        if conversation_id:
            raise ValueError(
                "Conversation documents cannot be uploaded "
                "inside the global document folder"
            )

    # --------------------------------------------------------
    # CONVERSATION VALIDATION
    # --------------------------------------------------------

    conversation_id = _validate_document_conversation(
        db=db,
        conversation_id=conversation_id,
        user_id=user_id,
    )

    # --------------------------------------------------------
    # CREATE DOCUMENT RECORD
    # --------------------------------------------------------

    doc = document_repo.create_file(
        db=db,
        file_name=file.filename,
        parent_id=parent_id,
        user_id=user_id,
        mime_type=file.content_type,
        conversation_id=conversation_id,
    )

    # --------------------------------------------------------
    # SAVE LOCAL UPLOAD
    # --------------------------------------------------------

    contents = file.file.read()

    file_path = os.path.join(
        UPLOAD_DIR,
        f"{doc.id}_{file.filename}",
    )

    with open(
        file_path,
        "wb",
    ) as f:
        f.write(contents)

    # --------------------------------------------------------
    # GCS PATH
    # --------------------------------------------------------

    gcs_path = generate_gcs_path(
        user_id=user_id,
        document_id=doc.id,
        original_filename=file.filename,
    )

    try:

        # ----------------------------------------------------
        # _process_for_rag is now a GENERATOR (it yields progress
        # events for the SSE upload flow). This non-streaming
        # variant has nowhere to send those events, but it still
        # MUST iterate the generator — otherwise the generator
        # object is created and immediately discarded, and none
        # of the actual extraction/embedding/indexing code inside
        # it ever runs. Draining it with a for-loop is what
        # actually executes the work.
        # ----------------------------------------------------

        for _ in _process_for_rag(
            db=db,
            doc=doc,
            file_path=file_path,
            conversation_id=conversation_id,
            user_id=user_id,
        ):
            pass

        doc = document_repo.update_file_storage_info(
            db=db,
            doc=doc,
            gcs_path=gcs_path,
            size_bytes=len(contents),
            status=DocumentStatus.READY,
        )

        logger.info(
            "Document upload and RAG processing completed: "
            "document_id=%s user_id=%s conversation_id=%s",
            doc.id,
            user_id,
            conversation_id,
        )

    except Exception:

        logger.exception(
            "RAG processing failed for document_id=%s "
            "user_id=%s conversation_id=%s",
            doc.id,
            user_id,
            conversation_id,
        )

        doc = document_repo.update_file_storage_info(
            db=db,
            doc=doc,
            gcs_path=gcs_path,
            size_bytes=len(contents),
            status=DocumentStatus.FAILED,
        )

    return doc


# ============================================================
# FILE TYPE DETECTION
# ============================================================


def _sniff_file_type(
    file_path: str,
) -> str:
    """
    Returns:
        image
        pdf
        unknown
    """

    try:

        with open(
            file_path,
            "rb",
        ) as f:
            header = f.read(8)

        if header.startswith(b"%PDF"):
            return "pdf"

    except Exception:

        return "unknown"

    try:

        with Image.open(file_path) as img:
            img.verify()

        return "image"

    except Exception:

        return "unknown"


def _detect_file_type(
    file_path: str,
    mime_type: Optional[str],
    file_name: str,
) -> str:
    """
    Determine RAG processing type.

    Returns:
        image
        pdf
        csv
        excel
        text
        docx
        pptx
        unknown
    """

    file_type = _sniff_file_type(
        file_path
    )

    if file_type != "unknown":
        return file_type

    if mime_type in IMAGE_MIME_TYPES:
        return "image"

    if mime_type == "application/pdf":
        return "pdf"

    if mime_type in DOCX_MIME_TYPES:
        return "docx"

    if mime_type in PPTX_MIME_TYPES:
        return "pptx"

    extension = os.path.splitext(
        file_name or ""
    )[1].lower()

    if extension in IMAGE_EXTENSIONS:
        return "image"

    if extension == ".pdf":
        return "pdf"

    if extension in CSV_EXTENSIONS:
        return "csv"

    if extension in EXCEL_EXTENSIONS:
        return "excel"

    if extension in DOCX_EXTENSIONS:
        return "docx"

    if extension in PPTX_EXTENSIONS:
        return "pptx"

    if extension in TEXT_EXTENSIONS:
        return "text"

    return "unknown"


# ============================================================
# RAG ROUTER
# ============================================================


def _process_for_rag(
    db: Session,
    doc: Document,
    file_path: str,
    conversation_id: Optional[str],
    user_id: str,
):
    """
    NOTE: this is now a GENERATOR. It yields progress dicts of the
    shape {"percent": int, "message": str} while it works, and the
    caller (upload_document_stream_service) forwards each one as
    its own SSE "processing" event so the frontend can render a
    live progress bar instead of jumping straight from 10% to 100%.
    """

    file_type = _detect_file_type(
        file_path=file_path,
        mime_type=doc.mime_type,
        file_name=doc.file_name,
    )

    yield {
        "percent": 15,
        "message": "Analyzing file...",
    }

    if file_type == "image":

        yield from _process_image_for_rag(
            db=db,
            doc=doc,
            file_path=file_path,
            conversation_id=conversation_id,
            user_id=user_id,
        )

        return

    if file_type == "pdf":

        yield from _process_pdf_for_rag(
            db=db,
            doc=doc,
            file_path=file_path,
            conversation_id=conversation_id,
            user_id=user_id,
        )

        return

    if file_type == "pptx":

        yield from _process_pptx_for_rag(
            db=db,
            doc=doc,
            file_path=file_path,
            conversation_id=conversation_id,
            user_id=user_id,
        )

        return

    if file_type in (
        "csv",
        "excel",
        "text",
        "docx",
    ):

        yield from _process_tabular_or_text_for_rag(
            db=db,
            doc=doc,
            file_path=file_path,
            file_type=file_type,
            conversation_id=conversation_id,
            user_id=user_id,
        )

        return

    # --------------------------------------------------------
    # UNKNOWN EXTENSION FALLBACK
    #
    # Any file type not explicitly recognised (.json, .yaml,
    # .log, .py, .html, .xml, source code, config files, etc.)
    # is still worth indexing if it's readable as plain text.
    # Only files that are genuinely binary/undecodable are
    # skipped. This is what lets the user upload "any type of
    # file" and still be able to query it. Genuinely unsupported
    # binary formats (.doc, .rtf) fall through to this branch
    # and will not produce searchable content — that is a real
    # gap, not a silent false-positive: no fabricated content is
    # ever indexed for them.
    # --------------------------------------------------------

    logger.info(
        "No specific handler for file type — attempting "
        "generic text fallback: document_id=%s mime_type=%s "
        "file_name=%s",
        doc.id,
        doc.mime_type,
        doc.file_name,
    )

    yield {
        "percent": 30,
        "message": "Extracting text...",
    }

    try:
        text = _extract_plain_text(file_path)

    except Exception:

        logger.exception(
            "Generic text fallback failed for document_id=%s "
            "— file is likely binary/unsupported.",
            doc.id,
        )

        return

    if not text or not text.strip():

        logger.warning(
            "Unsupported file type for RAG processing "
            "(no extractable text): document_id=%s "
            "mime_type=%s file_name=%s",
            doc.id,
            doc.mime_type,
            doc.file_name,
        )

        return

    yield {
        "percent": 60,
        "message": "Generating embeddings...",
    }

    chunk_count = _chunk_and_index_text(
        db=db,
        doc=doc,
        text=text,
        conversation_id=conversation_id,
        user_id=user_id,
    )

    yield {
        "percent": 95,
        "message": "Finalizing...",
    }

    logger.info(
        "RAG indexing completed via generic text fallback: "
        "document_id=%s text_chunks=%s user_id=%s "
        "conversation_id=%s",
        doc.id,
        chunk_count,
        user_id,
        conversation_id,
    )


# ============================================================
# STANDALONE IMAGE PROCESSING
# ============================================================


def _process_image_for_rag(
    db: Session,
    doc: Document,
    file_path: str,
    conversation_id: Optional[str],
    user_id: str,
):

    yield {
        "percent": 30,
        "message": "Analyzing image...",
    }

    try:

        caption = generate_image_caption(
            file_path
        )

    except ImageCaptionQuotaExceededError:

        logger.warning(
            "Image-captioning quota exhausted for "
            "document_id=%s — document uploaded but "
            "not caption-indexed yet.",
            doc.id,
        )

        return

    except Exception:

        logger.exception(
            "Image captioning failed for document_id=%s",
            doc.id,
        )

        return

    if not caption or not caption.strip():

        logger.warning(
            "Empty caption generated for document_id=%s",
            doc.id,
        )

        return

    yield {
        "percent": 65,
        "message": "Generating embeddings...",
    }

    chunks = [
        caption.strip()
    ]

    embeddings = embedding_manager.generate_embedding(
        chunks
    )

    yield {
        "percent": 90,
        "message": "Indexing...",
    }

    vector_store.add_documents(
        chunks=chunks,
        embeddings=embeddings,
        filename=doc.file_name,
        conversation_id=conversation_id,
        user_id=str(user_id),
        document_id=str(doc.id),
        content_type="image",
        parent_document_id=str(doc.id),
        page_number=None,
    )

    document_repo.create_chunks(
        db=db,
        doc_id=doc.id,
        chunks=chunks,
    )

    logger.info(
        "Image RAG indexing completed: "
        "document_id=%s user_id=%s conversation_id=%s",
        doc.id,
        user_id,
        conversation_id,
    )


# ============================================================
# SHARED TEXT CHUNKING + INDEXING
# ============================================================


def _chunk_and_index_text(
    db: Session,
    doc: Document,
    text: str,
    conversation_id: Optional[str],
    user_id: str,
) -> int:

    if not text or not text.strip():

        logger.warning(
            "No extractable text for document_id=%s",
            doc.id,
        )

        return 0

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        separators=[
            "\n\n",
            "\n",
            " ",
            "",
        ],
    )

    chunks = text_splitter.split_text(
        text
    )

    if not chunks:
        return 0

    embeddings = embedding_manager.generate_embedding(
        chunks
    )

    vector_store.add_documents(
        chunks=chunks,
        embeddings=embeddings,
        filename=doc.file_name,
        conversation_id=conversation_id,
        user_id=str(user_id),
        document_id=str(doc.id),
        content_type="text",
        parent_document_id=str(doc.id),
        page_number=None,
    )

    document_repo.create_chunks(
        db=db,
        doc_id=doc.id,
        chunks=chunks,
    )

    return len(chunks)


# ============================================================
# CSV
# ============================================================


def _extract_csv_text(
    file_path: str,
) -> str:

    lines = []

    with open(
        file_path,
        newline="",
        encoding="utf-8",
        errors="ignore",
    ) as f:

        reader = csv.reader(f)

        try:
            header = next(reader)

        except StopIteration:
            return ""

        for row in reader:

            row_text = ", ".join(
                f"{col}: {val}"
                for col, val in zip(
                    header,
                    row,
                )
            )

            lines.append(
                row_text
            )

    return "\n".join(lines)


# ============================================================
# EXCEL
# ============================================================


def _extract_excel_text(
    file_path: str,
) -> str:

    workbook = openpyxl.load_workbook(
        file_path,
        read_only=True,
        data_only=True,
    )

    lines = []

    try:

        for sheet in workbook.worksheets:

            rows = sheet.iter_rows(
                values_only=True
            )

            try:
                header = next(rows)

            except StopIteration:
                continue

            header = [
                str(h)
                if h is not None
                else ""
                for h in header
            ]

            lines.append(
                f"Sheet: {sheet.title}"
            )

            for row in rows:

                row_text = ", ".join(
                    f"{col}: {val}"
                    for col, val in zip(
                        header,
                        row,
                    )
                    if val is not None
                )

                if row_text:
                    lines.append(
                        row_text
                    )

    finally:

        workbook.close()

    return "\n".join(lines)


# ============================================================
# PLAIN TEXT
# ============================================================


def _extract_plain_text(
    file_path: str,
) -> str:

    with open(
        file_path,
        "r",
        encoding="utf-8",
        errors="ignore",
    ) as f:

        return f.read()


# ============================================================
# DOCX
# ============================================================


def _extract_docx_text(
    file_path: str,
) -> str:

    document = docx.Document(
        file_path
    )

    lines = []

    for paragraph in document.paragraphs:

        if paragraph.text.strip():

            lines.append(
                paragraph.text.strip()
            )

    for table in document.tables:

        for row in table.rows:

            row_text = ", ".join(
                cell.text.strip()
                for cell in row.cells
                if cell.text.strip()
            )

            if row_text:
                lines.append(
                    row_text
                )

    return "\n".join(lines)


# ============================================================
# CSV / EXCEL / TEXT / DOCX RAG
# ============================================================


def _process_tabular_or_text_for_rag(
    db: Session,
    doc: Document,
    file_path: str,
    file_type: str,
    conversation_id: Optional[str],
    user_id: str,
):

    yield {
        "percent": 30,
        "message": "Extracting text...",
    }

    try:

        if file_type == "csv":

            text = _extract_csv_text(
                file_path
            )

        elif file_type == "excel":

            text = _extract_excel_text(
                file_path
            )

        elif file_type == "docx":

            text = _extract_docx_text(
                file_path
            )

        else:

            text = _extract_plain_text(
                file_path
            )

    except Exception:

        logger.exception(
            "Content extraction failed for "
            "document_id=%s file_type=%s",
            doc.id,
            file_type,
        )

        return

    yield {
        "percent": 60,
        "message": "Generating embeddings...",
    }

    chunk_count = _chunk_and_index_text(
        db=db,
        doc=doc,
        text=text,
        conversation_id=conversation_id,
        user_id=user_id,
    )

    yield {
        "percent": 95,
        "message": "Finalizing...",
    }

    logger.info(
        "RAG indexing completed: document_id=%s "
        "file_type=%s text_chunks=%s "
        "user_id=%s conversation_id=%s",
        doc.id,
        file_type,
        chunk_count,
        user_id,
        conversation_id,
    )


# ============================================================
# DOCUMENT METADATA UPDATE
# ============================================================


def update_document_service(
    db: Session,
    doc: Document,
    file_name: Optional[str] = None,
    mime_type: Optional[str] = None,
) -> Document:

    if (
        file_name is None
        and mime_type is None
    ):
        return doc

    return document_repo.update_document(
        db=db,
        doc=doc,
        file_name=file_name,
        mime_type=mime_type,
    )


# ============================================================
# DELETE
# ============================================================


def delete_document_service(
    db: Session,
    doc: Document,
) -> bool:

    _delete_recursive(
        db=db,
        doc=doc,
    )

    return True


def _delete_recursive(
    db: Session,
    doc: Document,
) -> None:

    if doc.is_folder:

        children = document_repo.get_children(
            db,
            doc.id,
        )

        for child in children:

            _delete_recursive(
                db=db,
                doc=child,
            )

    else:

        # TODO:
        # GCS delete
        # gcs_client.delete(doc.gcs_path)

        # TODO:
        # Qdrant vector delete
        # vector_store.delete(document_id=doc.id)

        pass

    document_repo.delete_document_row(
        db,
        doc,
    )


# ============================================================
# PAGE-TAGGED IMAGE PERSISTENCE
#
# page_index is 0-based (as produced by PyMuPDF / python-pptx
# enumeration); the stored page_number is always the 1-based
# human-facing number ("page 4" / "slide 4").
# ============================================================


def _persist_page_image(
    db: Session,
    parent_doc: Document,
    image_bytes: bytes,
    extension: str,
    page_index: int,
    image_index: int,
    conversation_id: Optional[str],
    user_id: str,
) -> Document:

    extension = extension.lower()

    mime_ext = (
        "jpeg"
        if extension in (
            "jpg",
            "jpeg",
        )
        else extension
    )

    filename = (
        f"{parent_doc.file_name}"
        f"_p{page_index}"
        f"_{image_index}."
        f"{extension}"
    )

    image_doc = document_repo.create_file(
        db=db,
        file_name=filename,
        parent_id=parent_doc.id,
        user_id=user_id,
        mime_type=f"image/{mime_ext}",
        conversation_id=conversation_id,
        page_number=page_index + 1,
    )

    file_path = os.path.join(
        UPLOAD_DIR,
        f"{image_doc.id}_{filename}",
    )

    with open(
        file_path,
        "wb",
    ) as f:

        f.write(image_bytes)

    image_doc = document_repo.update_file_storage_info(
        db=db,
        doc=image_doc,
        gcs_path=generate_gcs_path(
            user_id,
            image_doc.id,
            filename,
        ),
        size_bytes=len(image_bytes),
        status=DocumentStatus.READY,
    )

    return image_doc


# Backward-compatible alias — some earlier internal callers may
# still reference the old name.
_persist_pdf_image = _persist_page_image


# ============================================================
# EXTRACT + PERSIST PDF IMAGES (embedded XObject images)
# ============================================================


def _extract_and_persist_pdf_images(
    db: Session,
    parent_doc: Document,
    file_path: str,
    conversation_id: Optional[str],
    user_id: str,
) -> list[tuple[Document, str, int]]:
    """
    Returns a list of (image_doc, image_path, page_index) tuples,
    one per persisted embedded image. page_index is 0-based and
    is used by the caller to know which pages already have at
    least one embedded image (to avoid re-rendering them in the
    per-page fallback below).
    """

    results = []

    pdf = fitz.open(
        file_path
    )

    try:

        for page_index in range(
            len(pdf)
        ):

            page = pdf[page_index]

            for image_index, img in enumerate(
                page.get_images(full=True)
            ):

                xref = img[0]

                try:

                    base_image = pdf.extract_image(
                        xref
                    )

                except Exception:

                    logger.warning(
                        "Unable to extract image "
                        "xref=%s page=%s "
                        "document_id=%s",
                        xref,
                        page_index,
                        parent_doc.id,
                    )

                    continue

                image_doc = _persist_page_image(
                    db=db,
                    parent_doc=parent_doc,
                    image_bytes=base_image["image"],
                    extension=base_image["ext"],
                    page_index=page_index,
                    image_index=image_index,
                    conversation_id=conversation_id,
                    user_id=user_id,
                )

                image_path = os.path.join(
                    UPLOAD_DIR,
                    f"{image_doc.id}_{image_doc.file_name}",
                )

                results.append(
                    (
                        image_doc,
                        image_path,
                        page_index,
                    )
                )

    finally:

        pdf.close()

    return results


# ============================================================
# SHARED CAPTION + OCR + EMBED + INDEX FOR A LIST OF IMAGES
#
# Used by both PDF and PPTX processing so the caption/OCR/embed
# logic exists in exactly one place (per your instruction not to
# duplicate functions when the same behavior is needed twice).
# ============================================================


def _caption_and_index_image_pairs(
    db: Session,
    parent_doc: Document,
    image_pairs: list[tuple[Document, str]],
    conversation_id: Optional[str],
    user_id: str,
    progress_start: int = 60,
    progress_end: int = 90,
):
    """
    image_pairs: list of (image_doc, image_path) — the page_index
    is not needed here since image_doc.page_number is already set.

    This is a GENERATOR: it yields a {"percent": ..., "message": ...}
    progress event after each image is processed (scaled between
    progress_start and progress_end, so the caller can place this
    loop wherever it falls in the overall pipeline), and RETURNS
    (image_count_indexed, image_count_skipped_due_to_limit) —
    callers should capture that via
    `result = yield from _caption_and_index_image_pairs(...)`.
    """

    image_count = 0
    skipped_count = 0

    images_to_process = image_pairs[
        :MAX_PDF_IMAGES_TO_CAPTION
    ]

    total_images = len(images_to_process)

    skipped_count = max(
        0,
        len(image_pairs)
        - MAX_PDF_IMAGES_TO_CAPTION,
    )

    if skipped_count:

        logger.warning(
            "Document has %s images; only captioning "
            "the first %s for document_id=%s. "
            "Remaining images are saved but "
            "not AI-caption indexed.",
            len(image_pairs),
            MAX_PDF_IMAGES_TO_CAPTION,
            parent_doc.id,
        )

    for image_index, (image_doc, image_path) in enumerate(
        images_to_process
    ):

        try:

            # ------------------------------------------------
            # GEMINI CAPTION
            # ------------------------------------------------

            caption = generate_image_caption(
                image_path
            )

            caption = (
                caption.strip()
                if caption
                else ""
            )

            # ------------------------------------------------
            # OCR
            # ------------------------------------------------

            ocr_text = ""

            try:

                with Image.open(
                    image_path
                ) as image:

                    ocr_text = (
                        pytesseract.image_to_string(
                            image
                        )
                        or ""
                    )

            except Exception:

                logger.warning(
                    "OCR failed for "
                    "image_document_id=%s",
                    image_doc.id,
                    exc_info=True,
                )

            ocr_text = ocr_text.strip()

            # ------------------------------------------------
            # COMBINED SEARCHABLE CONTENT
            # ------------------------------------------------

            searchable_parts = []

            if caption:

                searchable_parts.append(
                    "Image description:\n"
                    + caption
                )

            if ocr_text:

                searchable_parts.append(
                    "Text visible in image:\n"
                    + ocr_text
                )

            searchable_text = (
                "\n\n---\n\n".join(
                    searchable_parts
                )
            )

            if not searchable_text:

                logger.warning(
                    "No searchable content generated "
                    "for image_document_id=%s",
                    image_doc.id,
                )

                continue

            chunks = [
                searchable_text
            ]

            # ------------------------------------------------
            # EMBEDDING
            # ------------------------------------------------

            embeddings = (
                embedding_manager.generate_embedding(
                    chunks
                )
            )

            # ------------------------------------------------
            # QDRANT
            # ------------------------------------------------

            vector_store.add_documents(
                chunks=chunks,
                embeddings=embeddings,
                filename=image_doc.file_name,
                conversation_id=conversation_id,
                user_id=str(user_id),
                document_id=str(image_doc.id),
                content_type="image",
                parent_document_id=str(parent_doc.id),
                page_number=image_doc.page_number,
            )

            # ------------------------------------------------
            # DB CHUNK
            # ------------------------------------------------

            document_repo.create_chunks(
                db=db,
                doc_id=image_doc.id,
                chunks=chunks,
            )

            image_count += 1

            logger.info(
                "Image indexed successfully: "
                "image_document_id=%s "
                "parent_document_id=%s "
                "page_number=%s "
                "caption=%s "
                "ocr_chars=%s",
                image_doc.id,
                parent_doc.id,
                image_doc.page_number,
                bool(caption),
                len(ocr_text),
            )

        except ImageCaptionQuotaExceededError:

            remaining = (
                len(images_to_process)
                - image_count
            )

            skipped_count += max(
                0,
                remaining,
            )

            logger.warning(
                "Image-captioning quota exhausted "
                "while processing document_id=%s. "
                "Stopping image processing.",
                parent_doc.id,
            )

            break

        except Exception:

            logger.exception(
                "Captioning/indexing failed for "
                "image_document_id=%s "
                "parent_document_id=%s",
                image_doc.id,
                parent_doc.id,
            )

        # ----------------------------------------------------
        # PER-IMAGE PROGRESS
        #
        # Scaled between progress_start/progress_end so the
        # caller can place this loop anywhere in its own
        # percent range (e.g. 60-90 for a PDF that also had a
        # text phase before this).
        # ----------------------------------------------------

        current_percent = progress_start + int(
            (image_index + 1)
            / total_images
            * (progress_end - progress_start)
        )

        yield {
            "percent": min(
                current_percent,
                progress_end,
            ),
            "message": (
                f"Analyzing image {image_index + 1} "
                f"of {total_images}..."
            ),
        }

    return image_count, skipped_count


# ============================================================
# PDF RAG
# ============================================================


def _process_pdf_for_rag(
    db: Session,
    doc: Document,
    file_path: str,
    conversation_id: Optional[str],
    user_id: str,
):
    """
    PDF processing (GENERATOR — yields {"percent", "message"}
    progress events; see _process_for_rag's docstring):

    1. Normal PDF text (per page) -> chunks -> embeddings -> Qdrant

    2. Embedded PDF images ->
       child Document ->
       Gemini caption ->
       OCR ->
       combined searchable content ->
       embeddings -> Qdrant

    3. Per-page rendered fallback (§12): any page whose extracted
       text is very short (likely mostly a diagram/chart/
       screenshot) AND that produced zero embedded images is
       rendered and processed as an image in addition to whatever
       else that page has. This covers both the "fully scanned
       PDF" case (every page qualifies) and the narrower "one
       page has a vector-drawn chart" case that whole-document
       text/image emptiness checks miss.
    """

    # ========================================================
    # 1. NORMAL PDF TEXT — PER PAGE
    # ========================================================

    yield {
        "percent": 25,
        "message": "Extracting text from PDF...",
    }

    page_texts: list[str] = []

    try:

        with pdfplumber.open(
            file_path
        ) as pdf:

            for page in pdf.pages:

                page_text = (
                    page.extract_text()
                    or ""
                )

                page_texts.append(
                    page_text
                )

    except Exception:

        logger.exception(
            "PDF text extraction failed "
            "for document_id=%s",
            doc.id,
        )

        page_texts = []

    full_text = "\n".join(
        page_text
        for page_text in page_texts
        if page_text
    )

    text_chunks = []

    if full_text.strip():

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100,
            separators=[
                "\n\n",
                "\n",
                " ",
                "",
            ],
        )

        text_chunks = text_splitter.split_text(
            full_text
        )

    else:

        logger.warning(
            "No extractable text for document_id=%s",
            doc.id,
        )

    # ========================================================
    # 2. INDEX NORMAL PDF TEXT
    # ========================================================

    if text_chunks:

        yield {
            "percent": 40,
            "message": "Generating text embeddings...",
        }

        embeddings = (
            embedding_manager.generate_embedding(
                text_chunks
            )
        )

        vector_store.add_documents(
            chunks=text_chunks,
            embeddings=embeddings,
            filename=doc.file_name,
            conversation_id=conversation_id,
            user_id=str(user_id),
            document_id=str(doc.id),
            content_type="text",
            parent_document_id=str(doc.id),
            page_number=None,
        )

        document_repo.create_chunks(
            db=db,
            doc_id=doc.id,
            chunks=text_chunks,
        )

    # ========================================================
    # 3. EXTRACT EMBEDDED IMAGES
    # ========================================================

    yield {
        "percent": 50,
        "message": "Extracting images from PDF...",
    }

    embedded_results: list[tuple[Document, str, int]] = []

    try:

        embedded_results = (
            _extract_and_persist_pdf_images(
                db=db,
                parent_doc=doc,
                file_path=file_path,
                conversation_id=conversation_id,
                user_id=user_id,
            )
        )

    except Exception:

        logger.exception(
            "PDF image extraction failed "
            "for document_id=%s",
            doc.id,
        )

    pages_with_embedded_images = {
        page_index
        for _, _, page_index in embedded_results
    }

    image_pairs: list[tuple[Document, str]] = [
        (image_doc, image_path)
        for image_doc, image_path, _ in embedded_results
    ]

    # ========================================================
    # 4. PER-PAGE RENDERED FALLBACK (§12)
    #
    # Render any page that (a) has no embedded image of its own
    # and (b) has little/no extracted text — covers both fully
    # scanned PDFs and individual pages with a chart/diagram that
    # PyMuPDF's get_images() didn't expose as an XObject.
    # ========================================================

    total_pages = (
        len(page_texts)
        if page_texts
        else len(fitz.open(file_path))
    )

    pages_needing_render = [
        page_index
        for page_index in range(total_pages)
        if page_index not in pages_with_embedded_images
        and (
            page_index >= len(page_texts)
            or len(
                (page_texts[page_index] or "").strip()
            )
            < PDF_PAGE_TEXT_FALLBACK_THRESHOLD
        )
    ]

    if pages_needing_render:

        logger.info(
            "Rendering %s sparse-text/no-embedded-image page(s) "
            "as fallback images: document_id=%s pages=%s",
            len(pages_needing_render),
            doc.id,
            pages_needing_render,
        )

        yield {
            "percent": 55,
            "message": "Rendering scanned/visual pages...",
        }

        pdf = fitz.open(
            file_path
        )

        try:

            for page_index in pages_needing_render:

                page = pdf[page_index]

                pix = page.get_pixmap(
                    matrix=fitz.Matrix(
                        2,
                        2,
                    ),
                    alpha=False,
                )

                image_bytes = pix.tobytes(
                    "png"
                )

                image_doc = _persist_page_image(
                    db=db,
                    parent_doc=doc,
                    image_bytes=image_bytes,
                    extension="png",
                    page_index=page_index,
                    image_index=0,
                    conversation_id=conversation_id,
                    user_id=user_id,
                )

                image_path = os.path.join(
                    UPLOAD_DIR,
                    f"{image_doc.id}_{image_doc.file_name}",
                )

                image_pairs.append(
                    (
                        image_doc,
                        image_path,
                    )
                )

        finally:

            pdf.close()

    # ========================================================
    # 5. CAPTION + OCR + EMBED + INDEX ALL IMAGES
    # ========================================================

    if image_pairs:

        image_count, skipped_count = (
            yield from _caption_and_index_image_pairs(
                db=db,
                parent_doc=doc,
                image_pairs=image_pairs,
                conversation_id=conversation_id,
                user_id=user_id,
                progress_start=60,
                progress_end=95,
            )
        )

    else:

        image_count, skipped_count = 0, 0

        yield {
            "percent": 95,
            "message": "Finalizing...",
        }

    logger.info(
        "RAG indexing completed: "
        "document_id=%s "
        "text_chunks=%s "
        "image_chunks=%s "
        "image_chunks_skipped=%s "
        "user_id=%s "
        "conversation_id=%s",
        doc.id,
        len(text_chunks),
        image_count,
        skipped_count,
        user_id,
        conversation_id,
    )


# ============================================================
# PPTX TEXT EXTRACTION
# ============================================================


def _extract_pptx_text(
    file_path: str,
) -> str:

    presentation = Presentation(
        file_path
    )

    lines = []

    for slide_index, slide in enumerate(
        presentation.slides
    ):

        slide_lines = []

        for shape in slide.shapes:

            if (
                shape.has_text_frame
                and shape.text_frame.text.strip()
            ):

                slide_lines.append(
                    shape.text_frame.text.strip()
                )

            if shape.has_table:

                for row in shape.table.rows:

                    row_text = ", ".join(
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    )

                    if row_text:
                        slide_lines.append(
                            row_text
                        )

        if slide_lines:

            lines.append(
                f"Slide {slide_index + 1}:\n"
                + "\n".join(slide_lines)
            )

    return "\n\n".join(lines)


# ============================================================
# PPTX IMAGE EXTRACTION
# ============================================================


def _extract_and_persist_pptx_images(
    db: Session,
    parent_doc: Document,
    file_path: str,
    conversation_id: Optional[str],
    user_id: str,
) -> list[tuple[Document, str]]:

    presentation = Presentation(
        file_path
    )

    results = []

    # python-pptx picture shape type
    PICTURE_SHAPE_TYPE = 13

    for slide_index, slide in enumerate(
        presentation.slides
    ):

        image_index = 0

        for shape in slide.shapes:

            try:

                if shape.shape_type != PICTURE_SHAPE_TYPE:
                    continue

                image = shape.image

                image_doc = _persist_page_image(
                    db=db,
                    parent_doc=parent_doc,
                    image_bytes=image.blob,
                    extension=(
                        image.ext
                        or "png"
                    ),
                    page_index=slide_index,
                    image_index=image_index,
                    conversation_id=conversation_id,
                    user_id=user_id,
                )

                image_path = os.path.join(
                    UPLOAD_DIR,
                    f"{image_doc.id}_{image_doc.file_name}",
                )

                results.append(
                    (
                        image_doc,
                        image_path,
                    )
                )

                image_index += 1

            except Exception:

                logger.warning(
                    "Unable to extract PPTX image "
                    "slide=%s document_id=%s",
                    slide_index,
                    parent_doc.id,
                    exc_info=True,
                )

                continue

    return results


# ============================================================
# PPTX RAG
# ============================================================


def _process_pptx_for_rag(
    db: Session,
    doc: Document,
    file_path: str,
    conversation_id: Optional[str],
    user_id: str,
):
    """GENERATOR — yields {"percent", "message"} progress events."""

    # ========================================================
    # 1. TEXT
    # ========================================================

    yield {
        "percent": 25,
        "message": "Extracting text from slides...",
    }

    try:

        text = _extract_pptx_text(
            file_path
        )

    except Exception:

        logger.exception(
            "PPTX text extraction failed "
            "for document_id=%s",
            doc.id,
        )

        text = ""

    yield {
        "percent": 40,
        "message": "Generating text embeddings...",
    }

    chunk_count = _chunk_and_index_text(
        db=db,
        doc=doc,
        text=text,
        conversation_id=conversation_id,
        user_id=user_id,
    )

    # ========================================================
    # 2. IMAGES
    # ========================================================

    yield {
        "percent": 50,
        "message": "Extracting images from slides...",
    }

    image_pairs: list[tuple[Document, str]] = []

    try:

        image_pairs = (
            _extract_and_persist_pptx_images(
                db=db,
                parent_doc=doc,
                file_path=file_path,
                conversation_id=conversation_id,
                user_id=user_id,
            )
        )

    except Exception:

        logger.exception(
            "PPTX image extraction failed "
            "for document_id=%s",
            doc.id,
        )

    if image_pairs:

        image_count, skipped_count = (
            yield from _caption_and_index_image_pairs(
                db=db,
                parent_doc=doc,
                image_pairs=image_pairs,
                conversation_id=conversation_id,
                user_id=user_id,
                progress_start=60,
                progress_end=95,
            )
        )

    else:

        image_count, skipped_count = 0, 0

        yield {
            "percent": 95,
            "message": "Finalizing...",
        }

    logger.info(
        "RAG indexing completed: "
        "document_id=%s "
        "file_type=pptx "
        "text_chunks=%s "
        "image_chunks=%s "
        "image_chunks_skipped=%s "
        "user_id=%s "
        "conversation_id=%s",
        doc.id,
        chunk_count,
        image_count,
        skipped_count,
        user_id,
        conversation_id,
    )


# ============================================================
# SSE UPLOAD CONTEXT VALIDATION
# ============================================================


def resolve_upload_context(
    db: Session,
    parent_id: Optional[str],
    conversation_id: Optional[str],
    user_id: str,
):
    """
    Validate upload context BEFORE SSE starts.

    Returns:
        (parent_id, conversation_id)

    Returns None:
        parent folder does not exist.

    Raises:
        PermissionError
        ValueError
    """

    if not parent_id:
        parent_id = None

    if not conversation_id:
        conversation_id = None

    # --------------------------------------------------------
    # PARENT FOLDER VALIDATION
    # --------------------------------------------------------

    if parent_id:

        parent = document_repo.get_owned_folder_by_id(
            db,
            parent_id,
            user_id,
        )

        if not parent:
            return None

        # Conversation documents cannot be inside
        # global folders.
        if conversation_id:

            raise ValueError(
                "Conversation documents cannot be uploaded "
                "inside the global document folder"
            )

    # --------------------------------------------------------
    # CONVERSATION VALIDATION
    # --------------------------------------------------------

    conversation_id = (
        _validate_document_conversation(
            db=db,
            conversation_id=conversation_id,
            user_id=user_id,
        )
    )

    return (
        parent_id,
        conversation_id,
    )


# ============================================================
# SSE UPLOAD SERVICE
# ============================================================


def upload_document_stream_service(
    file: UploadFile,
    parent_id: Optional[str],
    conversation_id: Optional[str],
    user_id: str,
):
    """
    Generator used by the SSE upload endpoint.

    Events:

        uploading
        processing
        ready

    or:

        failed
    """

    db = SessionLocal()

    doc = None

    try:

        # ====================================================
        # UPLOADING
        # ====================================================

        yield {
            "status": DocumentStatus.UPLOADING.value,
            "percent": 5,
            "message": "Uploading file...",
        }

        # ====================================================
        # CREATE DOCUMENT
        # ====================================================

        doc = document_repo.create_file(
            db=db,
            file_name=file.filename,
            parent_id=parent_id,
            user_id=user_id,
            mime_type=file.content_type,
            conversation_id=conversation_id,
        )

        # ====================================================
        # SAVE FILE
        # ====================================================

        contents = file.file.read()

        file_path = os.path.join(
            UPLOAD_DIR,
            f"{doc.id}_{file.filename}",
        )

        with open(
            file_path,
            "wb",
        ) as f:

            f.write(contents)

        # ====================================================
        # GCS PATH
        # ====================================================

        gcs_path = generate_gcs_path(
            user_id=user_id,
            document_id=doc.id,
            original_filename=file.filename,
        )

        # ====================================================
        # PROCESSING STATUS
        # ====================================================

        doc = document_repo.update_document_status(
            db=db,
            doc=doc,
            status=DocumentStatus.PROCESSING,
        )

        yield {
            "status": DocumentStatus.PROCESSING.value,
            "document_id": str(doc.id),
            "file_name": doc.file_name,
            "percent": 10,
            "message": (
                "Processing document "
                "for knowledge base..."
            ),
        }

        # ====================================================
        # RAG PROCESSING
        # ====================================================

        try:

            # =================================================
            # RAG PROCESSING PROGRESS
            #
            # _process_for_rag is now a generator that yields
            # incremental {"percent": ...} progress events while
            # it works (text extraction, chunking, embedding,
            # image captioning, etc). Forward each one to the SSE
            # stream as its own "processing" event so the frontend
            # can render a live progress bar instead of jumping
            # straight from 10% to 100%.
            # =================================================

            for progress_event in _process_for_rag(
                db=db,
                doc=doc,
                file_path=file_path,
                conversation_id=conversation_id,
                user_id=user_id,
            ):

                yield {
                    "status": DocumentStatus.PROCESSING.value,
                    "document_id": str(doc.id),
                    "file_name": doc.file_name,
                    **progress_event,
                }

            # =================================================
            # READY
            # =================================================

            doc = document_repo.update_file_storage_info(
                db=db,
                doc=doc,
                gcs_path=gcs_path,
                size_bytes=len(contents),
                status=DocumentStatus.READY,
            )

            logger.info(
                "Document upload and RAG processing "
                "completed: document_id=%s "
                "user_id=%s conversation_id=%s",
                doc.id,
                user_id,
                conversation_id,
            )

            yield {
                "status": DocumentStatus.READY.value,
                "percent": 100,
                "document": (
                    DocumentOut.model_validate(
                        doc
                    ).model_dump(
                        mode="json"
                    )
                ),
                "message": (
                    "Document uploaded and "
                    "processed successfully"
                ),
            }

        except Exception:

            logger.exception(
                "Document processing failed: "
                "document_id=%s "
                "user_id=%s "
                "conversation_id=%s",
                doc.id,
                user_id,
                conversation_id,
            )

            # -----------------------------------------------
            # FAILED
            # -----------------------------------------------

            doc = document_repo.update_file_storage_info(
                db=db,
                doc=doc,
                gcs_path=gcs_path,
                size_bytes=len(contents),
                status=DocumentStatus.FAILED,
            )

            yield {
                "status": DocumentStatus.FAILED.value,
                "document": (
                    DocumentOut.model_validate(
                        doc
                    ).model_dump(
                        mode="json"
                    )
                ),
                "message": (
                    "Document processing failed"
                ),
            }

    except Exception:

        logger.exception(
            "Unexpected error during streamed upload "
            "for user_id=%s",
            user_id,
        )

        yield {
            "status": DocumentStatus.FAILED.value,
            "message": "Upload failed",
        }

    finally:

        db.close()